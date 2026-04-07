#!/usr/bin/env python3
"""Generate sample files for RAG system testing."""

import os
import pandas as pd
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches, Pt

os.makedirs("/home/ubuntu/sample_files", exist_ok=True)

# ===== 1. Create a Sample PDF =====
pdf_path = "/home/ubuntu/sample_files/sample_document.pdf"
doc = SimpleDocTemplate(pdf_path, pagesize=letter)
story = []
styles = getSampleStyleSheet()

title = Paragraph("Quantum Computing: A Comprehensive Guide", styles['Title'])
story.append(title)
story.append(Spacer(1, 0.3 * 72))

content = """
Quantum computing represents a fundamental shift in computational paradigm. Unlike classical computers that use bits (0 or 1), 
quantum computers leverage quantum bits or qubits, which can exist in a superposition of both states simultaneously.

Key Concepts:
1. Superposition: A qubit can be 0, 1, or both at the same time
2. Entanglement: Qubits can be correlated in ways that have no classical equivalent
3. Interference: Quantum algorithms manipulate probability amplitudes to increase correct answers and decrease wrong ones

Applications:
- Drug Discovery: Simulating molecular interactions
- Cryptography: Breaking RSA encryption
- Optimization: Solving complex combinatorial problems
- Machine Learning: Quantum machine learning algorithms
"""

for line in content.split('\n'):
    if line.strip():
        story.append(Paragraph(line, styles['Normal']))
        story.append(Spacer(1, 0.1 * 72))

story.append(Spacer(1, 0.2 * 72))
story.append(Paragraph("Quantum vs Classical Computers", styles['Heading2']))
story.append(Spacer(1, 0.1 * 72))

data = [
    ['Aspect', 'Classical', 'Quantum'],
    ['Basic Unit', 'Bit (0/1)', 'Qubit (0/1/both)'],
    ['Processing', 'Sequential', 'Parallel (Superposition)'],
    ['Speed', 'Exponential for some problems', 'Polynomial for some problems'],
    ['Error Rate', 'Low', 'High (Decoherence)'],
]

table = Table(data)
table.setStyle(TableStyle([
    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
    ('FONTSIZE', (0, 0), (-1, 0), 14),
    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
    ('GRID', (0, 0), (-1, -1), 1, colors.black),
]))
story.append(table)

doc.build(story)
print(f"✓ Created PDF: {pdf_path}")

# ===== 2. Create a Multi-Sheet Excel File =====
xlsx_path = "/home/ubuntu/sample_files/sample_data.xlsx"

with pd.ExcelWriter(xlsx_path, engine='openpyxl') as writer:
    # Sheet 1: Financial Data
    df1 = pd.DataFrame({
        'Quarter': ['Q1 2024', 'Q2 2024', 'Q3 2024', 'Q4 2024'],
        'Revenue': [1200000, 1450000, 1680000, 1920000],
        'Expenses': [800000, 920000, 1050000, 1200000],
        'Profit': [400000, 530000, 630000, 720000],
    })
    df1.to_excel(writer, sheet_name='Financial', index=False)
    
    # Sheet 2: Product Sales
    df2 = pd.DataFrame({
        'Product': ['Widget A', 'Widget B', 'Widget C', 'Widget D', 'Widget E'],
        'Units Sold': [5000, 7200, 3400, 8900, 6100],
        'Price per Unit': [25.50, 18.75, 42.00, 15.25, 31.99],
        'Total Revenue': [127500, 135000, 142800, 135725, 195099],
    })
    df2.to_excel(writer, sheet_name='Products', index=False)
    
    # Sheet 3: Customer Data
    df3 = pd.DataFrame({
        'Customer ID': ['C001', 'C002', 'C003', 'C004', 'C005'],
        'Name': ['Acme Corp', 'TechStart Inc', 'Global Solutions', 'Innovation Labs', 'Future Systems'],
        'Location': ['New York', 'San Francisco', 'London', 'Tokyo', 'Berlin'],
        'Annual Spend': [450000, 320000, 280000, 510000, 195000],
        'Account Status': ['Active', 'Active', 'Inactive', 'Active', 'Pending'],
    })
    df3.to_excel(writer, sheet_name='Customers', index=False)
    
    # Sheet 4: Market Analysis
    df4 = pd.DataFrame({
        'Region': ['North America', 'Europe', 'Asia Pacific', 'Latin America', 'Middle East'],
        'Market Size (M)': [2500, 1800, 3200, 600, 400],
        'Growth Rate (%)': [8.5, 5.2, 12.1, 6.8, 9.3],
        'Competitors': [12, 8, 15, 5, 3],
    })
    df4.to_excel(writer, sheet_name='Market', index=False)

print(f"✓ Created Multi-Sheet Excel: {xlsx_path}")

# ===== 3. Create a Sample PowerPoint =====
pptx_path = "/home/ubuntu/sample_files/sample_presentation.pptx"

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
left = Inches(1)
top = Inches(2.5)
width = Inches(8)
height = Inches(2)
txBox = slide1.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "Artificial Intelligence in Enterprise"
p = tf.paragraphs[0]
p.font.size = Pt(54)
p.font.bold = True

# Slide 2: Key Trends
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
tf = txBox.text_frame
tf.text = "Key AI Trends in 2024"
p = tf.paragraphs[0]
p.font.size = Pt(44)
p.font.bold = True

content_box = slide2.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(5))
tf = content_box.text_frame
tf.word_wrap = True
trends = [
    "• Generative AI adoption across industries",
    "• Large Language Models (LLMs) becoming mainstream",
    "• AI-powered automation in business processes",
    "• Ethical AI and responsible AI governance",
    "• Edge AI and on-device machine learning",
]
for i, trend in enumerate(trends):
    if i == 0:
        tf.text = trend
    else:
        p = tf.add_paragraph()
        p.text = trend
        p.level = 0
    tf.paragraphs[i].font.size = Pt(24)

# Slide 3: Implementation Strategy
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
tf = txBox.text_frame
tf.text = "Implementation Strategy"
p = tf.paragraphs[0]
p.font.size = Pt(44)
p.font.bold = True

strategy_box = slide3.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(5))
tf = strategy_box.text_frame
tf.word_wrap = True
strategies = [
    "Phase 1: Assessment & Planning (Q1-Q2)",
    "Phase 2: Pilot Projects (Q3)",
    "Phase 3: Scale & Optimize (Q4)",
    "Phase 4: Full Deployment (2025)",
]
for i, strategy in enumerate(strategies):
    if i == 0:
        tf.text = strategy
    else:
        p = tf.add_paragraph()
        p.text = strategy
        p.level = 0
    tf.paragraphs[i].font.size = Pt(22)

prs.save(pptx_path)
print(f"✓ Created PowerPoint: {pptx_path}")

# ===== 4. Create a Markdown File =====
md_path = "/home/ubuntu/sample_files/sample_guide.md"
md_content = """# Machine Learning Best Practices Guide

## Introduction
This guide covers essential best practices for implementing machine learning systems in production environments.

## 1. Data Preparation
- Ensure data quality and consistency
- Handle missing values appropriately
- Normalize and scale features
- Create balanced training/test sets

## 2. Model Selection
### Supervised Learning
- Linear Regression: For continuous predictions
- Classification Trees: For categorical predictions
- Neural Networks: For complex non-linear relationships

### Unsupervised Learning
- K-Means Clustering: For grouping similar data
- PCA: For dimensionality reduction
- Isolation Forest: For anomaly detection

## 3. Model Evaluation
| Metric | Use Case |
|--------|----------|
| Accuracy | Balanced datasets |
| Precision | When false positives are costly |
| Recall | When false negatives are costly |
| F1-Score | Balanced measure of precision and recall |

## 4. Deployment Considerations
1. **Monitoring**: Track model performance in production
2. **Versioning**: Maintain version control for models
3. **Rollback Plan**: Have a strategy to revert to previous models
4. **A/B Testing**: Compare new models against baselines

## 5. Common Pitfalls
- Data leakage from test to training sets
- Overfitting to training data
- Ignoring class imbalance
- Insufficient hyperparameter tuning
- Neglecting model interpretability

## Conclusion
Following these practices ensures robust, maintainable, and effective machine learning systems.
"""

with open(md_path, 'w') as f:
    f.write(md_content)
print(f"✓ Created Markdown: {md_path}")

print("\n✓ All sample files created successfully!")
print(f"Files location: /home/ubuntu/sample_files/")
